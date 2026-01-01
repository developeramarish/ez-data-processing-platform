#!/usr/bin/env python3
"""
EZ Platform Architecture Presentations Converter (Simple Version)
Simplified converter that works without Cairo dependencies.
"""

import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PresentationConverter:
    def __init__(self, base_dir):
        self.base_dir = Path(base_dir)
        self.output_dir = self.base_dir / "exports"
        self.output_dir.mkdir(exist_ok=True)

        # Create subdirectories for each format
        (self.output_dir / "pdf").mkdir(exist_ok=True)
        (self.output_dir / "markdown").mkdir(exist_ok=True)
        (self.output_dir / "powerpoint").mkdir(exist_ok=True)
        (self.output_dir / "svg").mkdir(exist_ok=True)

    def get_html_files(self, lang_dir=None):
        """Get all HTML files to convert."""
        if lang_dir:
            search_dir = self.base_dir / lang_dir
        else:
            search_dir = self.base_dir

        html_files = []
        for file in search_dir.glob("*.html"):
            if file.name != "index.html":  # Skip index
                html_files.append(file)
        return html_files

    def extract_content_from_html(self, html_path):
        """Extract all content from HTML file."""
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        # Extract title
        title = soup.find('h1')
        title_text = title.get_text() if title else html_path.stem

        # Extract subtitle
        subtitle = soup.find('p', class_='subtitle')
        subtitle_text = subtitle.get_text() if subtitle else ""

        # Extract SVG
        svg = soup.find('svg')
        svg_content = str(svg) if svg else None

        # Extract legend items
        legend_items = []
        legend = soup.find('div', class_='legend')
        if legend:
            for item in legend.find_all('div', class_='legend-item'):
                legend_items.append(item.get_text(strip=True))

        # Extract nav links
        nav_links = []
        nav = soup.find('div', class_='nav-links')
        if nav:
            for link in nav.find_all('a'):
                text = link.get_text(strip=True)
                href = link.get('href', '')
                nav_links.append({'text': text, 'href': href})

        return {
            'title': title_text,
            'subtitle': subtitle_text,
            'svg': svg_content,
            'legend': legend_items,
            'nav_links': nav_links,
            'soup': soup
        }

    def html_to_pdf(self, html_path, output_path):
        """Convert HTML to PDF using Playwright."""
        print(f"  [+] PDF: {output_path.name}")

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()

            # Load HTML file
            page.goto(f"file:///{html_path.absolute()}")

            # Wait for content to load
            page.wait_for_load_state('networkidle')

            # Generate PDF
            page.pdf(
                path=str(output_path),
                format='A4',
                landscape=True,
                print_background=True,
                margin={
                    'top': '0.5in',
                    'right': '0.5in',
                    'bottom': '0.5in',
                    'left': '0.5in'
                }
            )

            browser.close()

    def html_to_markdown_with_svg(self, html_path, output_path):
        """Convert HTML to Markdown with embedded SVG."""
        print(f"  [+] Markdown: {output_path.name}")

        content = self.extract_content_from_html(html_path)

        markdown = []
        markdown.append(f"# {content['title']}\n")

        if content['subtitle']:
            markdown.append(f"_{content['subtitle']}_\n")

        markdown.append("\n---\n\n")

        if content['svg']:
            # Write SVG to separate file
            svg_filename = output_path.stem + ".svg"
            svg_path = self.output_dir / "svg" / svg_filename

            with open(svg_path, 'w', encoding='utf-8') as f:
                f.write(content['svg'])

            print(f"  [+] SVG: {svg_filename}")

            # Embed SVG reference in markdown (relative path to svg directory)
            markdown.append(f"![{content['title']}](../svg/{svg_filename})\n\n")

            # Also embed inline SVG for direct viewing
            markdown.append("## Architecture Diagram (Embedded)\n\n")
            markdown.append(content['svg'])
            markdown.append("\n\n")

        # Add legend
        if content['legend']:
            markdown.append("## Legend\n\n")
            for item in content['legend']:
                markdown.append(f"- {item}\n")
            markdown.append("\n")

        # Add nav links
        if content['nav_links']:
            markdown.append("## Related Documentation\n\n")
            for link in content['nav_links']:
                text = link['text']
                href = link['href']
                if href and not href.startswith('http'):
                    # Convert HTML links to markdown links
                    md_href = href.replace('.html', '.md')
                    markdown.append(f"- [{text}](./{md_href})\n")
            markdown.append("\n")

        markdown.append("\n---\n\n")
        markdown.append(f"_Generated from {html_path.name}_\n")
        markdown.append(f"_Date: {Path(html_path).stat().st_mtime}_\n")

        # Write markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(''.join(markdown))

    def html_to_powerpoint(self, html_path, output_path):
        """Convert HTML to PowerPoint."""
        print(f"  [+] PowerPoint: {output_path.name}")

        # Extract content
        content = self.extract_content_from_html(html_path)

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Slide 1: Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title.text = content['title']
        subtitle_shape.text = content['subtitle']

        # Slide 2: Diagram reference (since we can't easily embed SVG)
        blank_slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title
        title.text = "Architecture Diagram"

        # Add text box with info
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = f"The architecture diagram is available in the following formats:\n\n"
        p.font.size = Pt(18)

        # Add bullet points
        formats = [
            f"• PDF: exports/pdf/{output_path.stem}.pdf",
            f"• SVG: exports/svg/{output_path.stem}.svg",
            f"• Markdown: exports/markdown/{output_path.stem}.md",
            f"• Original HTML: {html_path.name}"
        ]

        for fmt in formats:
            p = text_frame.add_paragraph()
            p.text = fmt
            p.font.size = Pt(16)
            p.space_before = Pt(6)

        # Slide 3: Legend (if exists)
        if content['legend']:
            slide = prs.slides.add_slide(blank_slide_layout)
            title = slide.shapes.title
            title.text = "Legend"

            left = Inches(1.5)
            top = Inches(2)
            width = Inches(10)
            height = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for item in content['legend']:
                p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.space_before = Pt(12)

        # Save PowerPoint
        prs.save(str(output_path))

    def convert_file(self, html_path, lang_suffix=""):
        """Convert a single HTML file to all formats."""
        print(f"\n[*] Converting: {html_path.name}")

        base_name = html_path.stem + lang_suffix

        # PDF
        pdf_output = self.output_dir / "pdf" / f"{base_name}.pdf"
        try:
            self.html_to_pdf(html_path, pdf_output)
        except Exception as e:
            print(f"  [X] ERROR creating PDF: {e}")

        # Markdown with SVG
        md_output = self.output_dir / "markdown" / f"{base_name}.md"
        try:
            self.html_to_markdown_with_svg(html_path, md_output)
        except Exception as e:
            print(f"  [X] ERROR creating Markdown: {e}")

        # PowerPoint
        pptx_output = self.output_dir / "powerpoint" / f"{base_name}.pptx"
        try:
            self.html_to_powerpoint(html_path, pptx_output)
        except Exception as e:
            print(f"  [X] ERROR creating PowerPoint: {e}")

    def convert_all(self):
        """Convert all HTML files in the directory."""
        print("=" * 80)
        print("EZ Platform Architecture Presentations Converter")
        print("=" * 80)

        # Convert English files
        print("\n[English Presentations]")
        en_files = self.get_html_files()
        for html_file in en_files:
            self.convert_file(html_file)

        # Convert Hebrew files
        print("\n[Hebrew Presentations]")
        he_files = self.get_html_files("he")
        for html_file in he_files:
            self.convert_file(html_file, lang_suffix="_he")

        print("\n" + "=" * 80)
        print("Conversion complete!")
        print("=" * 80)
        print(f"\nOutput directory: {self.output_dir.absolute()}\n")
        print(f"  PDF files:        {self.output_dir / 'pdf'}")
        print(f"  Markdown files:   {self.output_dir / 'markdown'}")
        print(f"  PowerPoint files: {self.output_dir / 'powerpoint'}")
        print(f"  SVG files:        {self.output_dir / 'svg'}")
        print("\n" + "=" * 80)

        # Count files
        pdf_count = len(list((self.output_dir / 'pdf').glob('*.pdf')))
        md_count = len(list((self.output_dir / 'markdown').glob('*.md')))
        pptx_count = len(list((self.output_dir / 'powerpoint').glob('*.pptx')))
        svg_count = len(list((self.output_dir / 'svg').glob('*.svg')))

        print(f"\nSummary:")
        print(f"  - {pdf_count} PDF files created")
        print(f"  - {md_count} Markdown files created")
        print(f"  - {pptx_count} PowerPoint files created")
        print(f"  - {svg_count} SVG diagrams extracted")
        print("\n" + "=" * 80)


def main():
    # Get the directory containing this script
    script_dir = Path(__file__).parent

    # Create converter and run
    converter = PresentationConverter(script_dir)
    converter.convert_all()


if __name__ == "__main__":
    main()
